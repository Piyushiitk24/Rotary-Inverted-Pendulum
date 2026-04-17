from __future__ import annotations

import copy
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path


def repo_root() -> Path:
    return Path(__file__).resolve().parents[3]


def workspace_root() -> Path:
    return Path(__file__).resolve().parents[1]


def scripts_root() -> Path:
    return Path(__file__).resolve().parent


def load_manifest() -> dict:
    sys.path.insert(0, str(scripts_root()))
    from content_manifest import get_manifest

    return get_manifest()


def prepare_dirs(root: Path) -> None:
    for rel in [
        "assets",
        "assets/equations",
        "assets/figures",
        "dist",
        ".mpl-cache",
    ]:
        (root / rel).mkdir(parents=True, exist_ok=True)


def clean_dir(path: Path) -> None:
    if not path.exists():
        return
    for child in path.iterdir():
        if child.is_dir():
            shutil.rmtree(child)
        else:
            child.unlink()


def clean_generated(workspace: Path) -> None:
    for rel in ["assets/equations", "assets/figures", "dist"]:
        clean_dir(workspace / rel)


def copy_figures(manifest: dict, repo: Path, workspace: Path) -> dict:
    manifest_copy = copy.deepcopy(manifest)
    seen: dict[str, str] = {}
    for slide in manifest_copy["slides"]:
        for element in slide.get("elements", []):
            if element.get("type") != "image":
                continue
            src = repo / element["src"]
            if not src.exists():
                raise FileNotFoundError(f"Missing image source: {src}")
            ext = src.suffix.lower()
            asset_rel = Path("assets") / "figures" / f"{element['asset_id']}{ext}"
            dst = workspace / asset_rel
            if element["asset_id"] not in seen:
                shutil.copy2(src, dst)
                seen[element["asset_id"]] = str(asset_rel)
            element["asset_path"] = seen[element["asset_id"]]
    return manifest_copy


def write_manifest(manifest: dict, path: Path) -> None:
    path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")


def run(cmd: list[str], *, cwd: Path, env: dict[str, str] | None = None) -> None:
    subprocess.run(cmd, cwd=str(cwd), env=env, check=True)


def verify_deck(deck_path: Path, expected_slides: int) -> None:
    from pptx import Presentation

    prs = Presentation(str(deck_path))
    if len(prs.slides) != expected_slides:
        raise RuntimeError(f"Expected {expected_slides} slides, found {len(prs.slides)}")


def main() -> int:
    workspace = workspace_root()
    repo = repo_root()
    clean_generated(workspace)
    prepare_dirs(workspace)

    manifest = load_manifest()
    if len(manifest["slides"]) != 21:
        raise RuntimeError(f"Expected 21 slides in manifest, found {len(manifest['slides'])}")
    manifest = copy_figures(manifest, repo, workspace)
    manifest_path = workspace / "dist" / manifest["meta"]["manifest_filename"]
    write_manifest(manifest, manifest_path)

    env = os.environ.copy()
    env["MPLCONFIGDIR"] = str(workspace / ".mpl-cache")
    env["TEXINPUTS"] = str(workspace / "tex") + os.pathsep

    base_deck = workspace / "dist" / manifest["meta"]["base_deck_filename"]
    final_deck = workspace / "dist" / manifest["meta"]["deck_filename"]

    run(
        [sys.executable, str(scripts_root() / "render_equations.py"), str(manifest_path), str(workspace)],
        cwd=workspace,
        env=env,
    )
    run(
        ["node", str(scripts_root() / "build_deck.js"), str(manifest_path), str(base_deck)],
        cwd=workspace,
        env=env,
    )
    run(
        [sys.executable, str(scripts_root() / "postprocess_deck.py"), str(manifest_path), str(workspace), str(base_deck), str(final_deck)],
        cwd=workspace,
        env=env,
    )

    verify_deck(final_deck, 21)
    print(f"[OK] Deck written to {final_deck}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
