from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def load_manifest(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def fit_contain(path: Path, box_w: float, box_h: float) -> tuple[float, float]:
    with Image.open(path) as img:
        width, height = img.size
    aspect = width / height
    box_aspect = box_w / box_h
    if aspect >= box_aspect:
        return box_w, box_w / aspect
    return box_h * aspect, box_h


def add_picture(slide, img_path: Path, x: float, y: float, w: float, h: float) -> None:
    fit_w, fit_h = fit_contain(img_path, w, h)
    left = x + (w - fit_w) / 2.0
    top = y + (h - fit_h) / 2.0
    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), Inches(fit_w), Inches(fit_h))


def main() -> int:
    if len(sys.argv) != 5:
        print("Usage: python postprocess_deck.py <manifest.json> <workspace_root> <base_deck.pptx> <final_deck.pptx>", file=sys.stderr)
        return 2

    manifest_path = Path(sys.argv[1]).resolve()
    workspace_root = Path(sys.argv[2]).resolve()
    base_deck = Path(sys.argv[3]).resolve()
    final_deck = Path(sys.argv[4]).resolve()

    manifest = load_manifest(manifest_path)
    prs = Presentation(str(base_deck))

    for slide_def, slide in zip(manifest["slides"], prs.slides, strict=True):
        for element in slide_def.get("elements", []):
            if element.get("type") == "image":
                img_path = workspace_root / element["asset_path"]
                add_picture(slide, img_path, element["x"], element["y"], element["w"], element["h"])
            elif element.get("type") == "equation":
                img_path = workspace_root / "assets" / "equations" / f"{element['eq_id']}.png"
                add_picture(slide, img_path, element["x"], element["y"], element["w"], element["h"])

    prs.save(str(final_deck))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
